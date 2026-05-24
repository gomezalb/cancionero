"""
generar_canciones.py
====================
Recorre tu cancionero (carpetas A/, B/, C/...) y genera canciones.json
listo para usar en la app web del cancionero.

USO
---
1. Instalá la dependencia (una sola vez):
       pip install python-pptx

2. Editá la variable CARPETA_CANCIONERO con tu ruta real de OneDrive.

3. Corré el script:
       python generar_canciones.py

4. Se genera canciones.json en la misma carpeta del script.
   Subís ese archivo a GitHub junto con el index.html.

MARCAR CANCIONES COMO NUEVAS
-----------------------------
Renombrá el archivo agregando .NEW antes de .pptx:

    Nombre de la cancion.NEW.pptx

El script detecta el sufijo, guarda "nueva": true en el JSON
y limpia el .NEW del título automáticamente.

FORMATO PRESERVADO
------------------
  - Saltos de línea: respetados tal como están en el .pptx
  - Negrita: <b>...</b>
  - Contadores tipo "1/1", "2/3" se filtran automáticamente

DIAGNÓSTICO
-----------
    python generar_canciones.py --debug
"""

import json
import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn

# ─────────────────────────────────────────────
# CONFIGURACIÓN — editá estas líneas
# ─────────────────────────────────────────────
CARPETA_CANCIONERO = r"C:\Users\alber\OneDrive\Documentos\Cancionero"
ARCHIVO_SALIDA     = "canciones.json"
# ─────────────────────────────────────────────

DEBUG = "--debug" in sys.argv
PATRON_CONTADOR = re.compile(r"^\s*\d+\s*/\s*\d+\s*$")


def html_escape(text):
    return (text
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;"))


def get_run_bold(r_elem):
    try:
        rPr = r_elem.find(qn("a:rPr"))
        if rPr is not None:
            return rPr.get("b") in ("1", "true")
    except Exception:
        pass
    return False


def run_elem_a_html(r_elem):
    t_elem = r_elem.find(qn("a:t"))
    if t_elem is None or not t_elem.text:
        return ""
    texto = t_elem.text
    if PATRON_CONTADOR.match(texto):
        return ""
    texto_esc = html_escape(texto)
    if get_run_bold(r_elem):
        texto_esc = f"<b>{texto_esc}</b>"
    return texto_esc


def parrafo_a_html(p_elem):
    partes = []
    for child in p_elem:
        tag = child.tag
        if tag == qn("a:r"):
            html = run_elem_a_html(child)
            if html:
                partes.append(html)
        elif tag == qn("a:br"):
            partes.append("<br>")

    contenido = "".join(partes).strip()
    texto_plano = "".join(
        (child.find(qn("a:t")).text or "")
        for child in p_elem.findall(qn("a:r"))
        if child.find(qn("a:t")) is not None
    ).strip()

    if PATRON_CONTADOR.match(texto_plano):
        return None
    return contenido if contenido else None


def extraer_html_slide(slide):
    lineas = []
    for shape in slide.shapes:
        try:
            if not shape.has_text_frame:
                continue
            txBody = shape.text_frame._txBody
            for p_elem in txBody.findall(qn("a:p")):
                html_para = parrafo_a_html(p_elem)
                if html_para is None:
                    if lineas and lineas[-1] != "":
                        lineas.append("")
                else:
                    lineas.append(html_para)
        except Exception:
            continue

    while lineas and lineas[0] == "":
        lineas.pop(0)
    while lineas and lineas[-1] == "":
        lineas.pop()

    if not lineas:
        return ""

    partes = []
    for linea in lineas:
        if linea == "":
            partes.append("<br>")
        else:
            partes.append(linea)
    html = "<br>".join(partes)

    # Si toda la slide está en negrita, reemplazar múltiples <b>...</b>
    # por un único par que envuelva todo el contenido
    sin_bold = html.replace("<b>", "").replace("</b>", "")
    if sin_bold != html:
        # Había tags de negrita → envolver todo en uno solo
        html = f"<b>{sin_bold}</b>"

    return html


def extraer_texto_plano_slide(slide):
    textos = []
    for shape in slide.shapes:
        try:
            if not shape.has_text_frame:
                continue
            texto = shape.text_frame.text.strip()
            if texto and not PATRON_CONTADOR.match(texto):
                textos.append(texto)
        except Exception:
            continue
    return " ".join(textos)


def primera_slide_es_solo_titulo(slide, titulo):
    textos = []
    for shape in slide.shapes:
        try:
            if not shape.has_text_frame:
                continue
            texto = shape.text_frame.text.strip()
            if texto and not PATRON_CONTADOR.match(texto):
                textos.append(texto)
        except Exception:
            continue
    return "\n".join(textos).strip() == titulo.strip()


def procesar_pptx(ruta_archivo):
    """
    Devuelve (resultado_dict, motivo_omision).

    Sufijos soportados en el nombre del archivo (en cualquier orden):
      .NEW       → canción nueva
      .DO .REm   → tonalidad (cualquier nota con o sin 'm' para menor)
      .FAs .FAsm → sostenido: usá 's' como alias de '#' (FA#, FA#m, etc.)
      .FA#m          → también se acepta el símbolo # directamente

    Ejemplos:
      Bésame mucho.LAm.pptx
      Gracias a la vida.NEW.SOL.pptx
      Caminito.RE.NEW.pptx
      Aleluya.FAs.pptx          → guarda tono "FA#"
      Aleluya.FAsm.NEW.pptx     → guarda tono "FA#m", marcada como nueva
    """
    stem = ruta_archivo.stem  # nombre sin extensión

    # Detectar y extraer sufijos .NEW y tonalidad en cualquier orden
    # Tonalidades válidas: DO RE MI FA SOL LA SI
    # + variantes con 'm' (menor) y 's' (sostenido, alias de #)
    # Ejemplos: .FA .FAm .FAs .FAsm (FA# menor)
    PATRON_TONO = re.compile(
        r"\.(DO|RE|MI|FA|SOL|LA|SI)([s#]?)(m?)$",
        re.IGNORECASE
    )

    es_nueva = False
    tono = None
    partes = stem

    # Procesar sufijos iterativamente (puede haber varios)
    for _ in range(4):
        if partes.upper().endswith(".NEW"):
            es_nueva = True
            partes = partes[:-4]
        elif m := PATRON_TONO.search(partes):
            nota      = m.group(1).upper()
            sostenido = m.group(2).lower()   # 's' o ''
            menor     = m.group(3).lower()   # 'm' o ''
            # Normalizar: 's' -> '#' para mostrar bien en pantalla
            tono = nota + ("#" if sostenido in ("s", "#") else "") + menor
            partes = partes[:m.start()]
        else:
            break

    titulo = partes.strip()

    try:
        prs = Presentation(ruta_archivo)
    except Exception as e:
        return None, f"no se pudo abrir: {e}"

    if not prs.slides:
        return None, "no tiene diapositivas"

    slides_list = list(prs.slides)
    primera = slides_list[0]
    preview = extraer_texto_plano_slide(primera)

    slides_html = []
    for i, slide in enumerate(slides_list):
        if i == 0 and primera_slide_es_solo_titulo(slide, titulo):
            continue
        html = extraer_html_slide(slide)
        if html:
            slides_html.append(html)

    if not slides_html:
        return None, "todas las slides quedaron vacías"

    resultado = {
        "titulo": titulo,
        "preview": preview,
        "slides": slides_html
    }
    if es_nueva:
        resultado["nueva"] = True
    if tono:
        resultado["tono"] = tono

    return resultado, None


def generar_json(carpeta_raiz, archivo_salida):
    raiz = Path(carpeta_raiz)

    if not raiz.exists():
        print(f"❌ La carpeta no existe: {raiz}")
        return

    carpetas = sorted([d for d in raiz.iterdir() if d.is_dir()])
    if not carpetas:
        print(f"❌ No se encontraron subcarpetas en: {raiz}")
        return

    # Cargar acordes existentes del JSON anterior para preservarlos
    acordes_existentes = {}
    salida = Path(archivo_salida)
    if salida.exists():
        try:
            with open(salida, encoding="utf-8") as f:
                anterior = json.load(f)
            for c in anterior:
                if c.get("acordes"):
                    acordes_existentes[c["titulo"]] = c["acordes"]
            if acordes_existentes:
                print(f"💾 Preservando acordes de {len(acordes_existentes)} canción(es)\n")
        except Exception:
            pass

    print(f"📂 Cancionero: {raiz}")
    print(f"📁 Carpetas: {[d.name for d in carpetas]}\n")

    canciones = []
    omitidas = []

    for carpeta in carpetas:
        nombre_carpeta = carpeta.name
        archivos = sorted(carpeta.glob("*.pptx"))

        if not archivos:
            print(f"  [{nombre_carpeta}] — sin .pptx, se omite")
            continue

        print(f"  [{nombre_carpeta}] — {len(archivos)} archivo(s)")

        for archivo in archivos:
            resultado, motivo = procesar_pptx(archivo)

            if resultado is None:
                print(f"    ✗ {archivo.name} — {motivo}")
                omitidas.append((str(archivo), motivo))
            else:
                resultado["carpeta"] = nombre_carpeta
                # Preservar acordes si existen en el JSON anterior
                if resultado["titulo"] in acordes_existentes:
                    resultado["acordes"] = acordes_existentes[resultado["titulo"]]
                canciones.append(resultado)
                diap = len(resultado["slides"])
                nueva_tag = " [NUEVA]" if resultado.get("nueva") else ""
                tono_tag = f" [{resultado['tono']}]" if resultado.get("tono") else ""
                acorde_tag = " 🎵" if resultado.get("acordes") else ""
                print(f"    ✓ {resultado['titulo']:<45}{nueva_tag}{tono_tag}{acorde_tag} ({diap} diapositiva{'s' if diap != 1 else ''})")

    canciones.sort(key=lambda c: c["titulo"].lower())

    with open(salida, "w", encoding="utf-8") as f:
        json.dump(canciones, f, ensure_ascii=False, indent=2)

    nuevas = sum(1 for c in canciones if c.get("nueva"))
    con_acordes = sum(1 for c in canciones if c.get("acordes"))
    print(f"\n✅ {len(canciones)} canciones guardadas en: {salida.resolve()}")
    if nuevas:
        print(f"   {nuevas} marcadas como nuevas")
    if con_acordes:
        print(f"   {con_acordes} con acordes preservados")
    if omitidas:
        print(f"   {len(omitidas)} archivos omitidos")
        if not DEBUG:
            print("   (corré con --debug para ver el detalle)")


if __name__ == "__main__":
    generar_json(CARPETA_CANCIONERO, ARCHIVO_SALIDA)
